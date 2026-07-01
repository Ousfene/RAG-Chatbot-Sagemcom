"""
Professional Text Extraction Module (Per-Page/Section)
-----------------------------------------------------
- Extracts text per page (PDF), per section (DOCX), per slide (PPTX), per sheet (XLSX), or as a whole (TXT).
- Each chunk includes file_name, file_path, file_type, page/section/slide, and content.
"""
import fitz
import pdfplumber
import docx
import pptx
import openpyxl
from pathlib import Path
from typing import Dict, List, Any
import json
import logging
import re
OUTPUT_DIR = Path("./data_prep/processed/texts")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

logging.basicConfig(level=logging.INFO, format='%(asctime)s [%(levelname)s] %(message)s')
logger = logging.getLogger("extract_texts")

def is_heading(line, font_size, font_name):
    line = line.strip()
    # Only treat as heading if font is large or bold and line is not too short/long and not a date/author line
    if len(line) < 6 or len(line) > 80:
        return False
    # Exclude lines with dates or author patterns
    if re.search(r"\d{2}\.\d{2}\.\d{4}", line) or "|" in line or "Ajout" in line or "Création" in line:
        return False
    # Exclude lines with long runs of dots (likely table of contents or table rows)
    if re.search(r"\.{5,}", line):
        return False
    if font_size >= 13:
        return True
    if ("Bold" in font_name or "bold" in font_name):
        return True
    # Numbered section pattern (e.g., "4.1 Abréviations", "3.2 Logigramme") - must be number, dot, number, space, then text
    if re.match(r"^\d+(\.\d+)*\s+.+", line):
        return True
    # All caps, not too short, and not just a single word
    if line.isupper() and " " in line:
        return True
    return False

def extract_text_from_pdf(pdf_path: Path) -> List[Dict[str, Any]]:
    """Extract structured blocks from a PDF using headings and font info."""
    results = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                words = page.extract_words(extra_attrs=["fontname", "size"])
                # Group words by line
                lines = []
                current_line = []
                last_top = None
                for word in words:
                    if "top" not in word:
                        continue
                    top = word["top"]
                    if last_top is not None and abs(top - last_top) > 2:
                        if current_line:
                            lines.append(current_line)
                        current_line = []
                    current_line.append(word)
                    last_top = top
                if current_line:
                    lines.append(current_line)
                # --- DEBUG: Print all lines for page 5 ---
                if page_num == 5:
                    print("Extracted lines for page 5:")
                    for line in lines:
                        print(" ".join(w["text"] for w in line))
                # Convert lines to text and font info
                line_infos = []
                for line in lines:
                    if not line:
                        continue
                    text = " ".join(w["text"] for w in line).strip()
                    font_sizes = [w.get("size", 10) for w in line]
                    font_names = [w.get("fontname", "") for w in line]
                    font_size = max(font_sizes) if font_sizes else 10
                    font_name = font_names[0] if font_names else ""
                    line_infos.append((text, font_size, font_name))
                # Detect blocks with parent_heading support (improved for consecutive headings)
                blocks = []
                current_block = {"heading": None, "content": [], "page": page_num, "parent_heading": None}
                last_heading = None
                for text, font_size, font_name in line_infos:
                    if is_heading(text, font_size, font_name):
                        # Always save previous block (even if empty, for parent_heading)
                        if current_block["heading"] is not None or current_block["content"]:
                            blocks.append(current_block)
                            if not current_block["content"]:
                                last_heading = current_block["heading"]
                        # New block: set parent_heading if last_heading exists and not same as this heading
                        parent_heading = last_heading if last_heading and last_heading != text else None
                        current_block = {"heading": text, "content": [], "page": page_num, "parent_heading": parent_heading}
                    else:
                        current_block["content"].append(text)
                # Always save the last block
                if current_block["heading"] is not None or current_block["content"]:
                    blocks.append(current_block)
                # Add metadata
                for block in blocks:
                    block["file_name"] = pdf_path.name
                    block["file_path"] = str(pdf_path)
                    block["file_type"] = ".pdf"
                    block["source_link"] = f"file:///{pdf_path.resolve().as_posix()}#page={page_num}"
                    # Keep content as a list of lines for structure (not joined with \n)
                    results.append(block)
    except Exception as ex:
        logger.error(f"[!!] PDF extraction failed for {pdf_path.name}: {ex}")
        # Fallback to PyMuPDF (page-level only)
        with fitz.open(pdf_path) as doc:
            for i, page in enumerate(doc, 1):
                text = page.get_text("text")
                results.append({
                    "file_name": pdf_path.name,
                    "file_path": str(pdf_path),
                    "file_type": ".pdf",
                    "page": i,
                    "content": text
                })
    return results

def extract_text_from_docx(docx_path: Path) -> List[Dict[str, Any]]:
    """Extract text per paragraph from DOCX, grouped by section (heading if available)."""
    doc = docx.Document(docx_path)
    results = []
    section_num = 0
    section_text = []
    section_title = None
    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            # Save previous section
            if section_text:
                results.append({
                    "file_name": docx_path.name,
                    "file_path": str(docx_path),
                    "file_type": ".docx",
                    "section": section_num,
                    "section_title": section_title,
                    "content": "\n".join(section_text)
                })
                section_text = []
            section_num += 1
            section_title = para.text.strip()
        else:
            section_text.append(para.text)
    # Save last section
    if section_text:
        results.append({
            "file_name": docx_path.name,
            "file_path": str(docx_path),
            "file_type": ".docx",
            "section": section_num if section_num else 1,
            "section_title": section_title,
            "content": "\n".join(section_text)
        })
    return results

def extract_text_from_pptx(pptx_path: Path) -> List[Dict[str, Any]]:
    """Extract text per slide from PPTX."""
    prs = pptx.Presentation(pptx_path)
    results = []
    for i, slide in enumerate(prs.slides, 1):
        slide_text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        results.append({
            "file_name": pptx_path.name,
            "file_path": str(pptx_path),
            "file_type": ".pptx",
            "slide": i,
            "content": slide_text
        })
    return results

def extract_text_from_xlsx(xlsx_path: Path) -> List[Dict[str, Any]]:
    """Extract text per sheet from XLSX."""
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    results = []
    for i, sheet in enumerate(wb.worksheets, 1):
        texts = []
        for row in sheet.iter_rows():
            texts.append("\t".join(str(cell.value or "") for cell in row))
        results.append({
            "file_name": xlsx_path.name,
            "file_path": str(xlsx_path),
            "file_type": ".xlsx",
            "sheet": i,
            "sheet_title": sheet.title,
            "content": "\n".join(texts)
        })
    return results

def extract_text_from_txt(txt_path: Path) -> List[Dict[str, Any]]:
    """Extract all text from TXT as a single chunk."""
    text = txt_path.read_text(encoding="utf-8", errors="ignore")
    return [{
        "file_name": txt_path.name,
        "file_path": str(txt_path),
        "file_type": ".txt",
        "content": text
    }]

EXTRACTORS = {
    ".pdf": extract_text_from_pdf,
    ".docx": extract_text_from_docx,
    ".pptx": extract_text_from_pptx,
    ".xlsx": extract_text_from_xlsx,
    ".txt": extract_text_from_txt,
}

def extract_all_texts(raw_dir: Path = Path("raw_files"), output_dir: Path = OUTPUT_DIR):
    output_dir.mkdir(parents=True, exist_ok=True)
    for file_path in raw_dir.rglob("*"):
        if not file_path.is_file():
            continue
        ext = file_path.suffix.lower()
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            continue
        logger.info(f"[text] Extracting {file_path.name}")
        try:
            chunks = extractor(file_path)
        except Exception as ex:
            logger.error(f"[!!] Failed while extracting text from {file_path.name}: {ex}")
            continue
        try:
            output_file = output_dir / (file_path.stem + ".json")
            with output_file.open("w", encoding="utf-8") as f:
                json.dump(chunks, f, ensure_ascii=False, indent=2)
            logger.info(f"[text] Wrote {len(chunks)} chunks to {output_file}")
        except Exception as e:
            logger.error(f"[!!] Failed to write JSON for {file_path.name}: {e}")

def extract_all_texts_with_summary(raw_dir: Path = Path("raw_files"), output_dir: Path = OUTPUT_DIR):
    output_dir.mkdir(parents=True, exist_ok=True)
    summary = {"files": 0, "chunks": 0, "errors": 0}
    for file_path in raw_dir.rglob("*"):
        if not file_path.is_file():
            continue
        ext = file_path.suffix.lower()
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            continue
        logger.info(f"[text] Extracting {file_path.name}")
        try:
            chunks = extractor(file_path)
            summary["files"] += 1
            summary["chunks"] += len(chunks)
        except Exception as ex:
            logger.error(f"[!!] Failed while extracting text from {file_path.name}: {ex}")
            summary["errors"] += 1
            continue
        try:
            output_file = output_dir / (file_path.stem + ".json")
            with output_file.open("w", encoding="utf-8") as f:
                json.dump(chunks, f, ensure_ascii=False, indent=2)
            logger.info(f"[text] Wrote {len(chunks)} chunks to {output_file}")
        except Exception as e:
            logger.error(f"[!!] Failed to write JSON for {file_path.name}: {e}")
            summary["errors"] += 1
    logger.info(f"[summary] Files processed: {summary['files']}, Chunks: {summary['chunks']}, Errors: {summary['errors']}")

if __name__ == "__main__":
    extract_all_texts_with_summary()
