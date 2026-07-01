# modules/extract_images.py

import os
import pathlib
import fitz  # PyMuPDF for PDF image extraction
import docx
import pptx
import openpyxl
from pathlib import Path
from typing import Dict
import shutil
import uuid

OUTPUT_DIR = Path("./data_prep/processed/images")
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def save_image(image_bytes: bytes, dest_folder: Path, file_stem: str, image_id: str):
    dest_folder.mkdir(parents=True, exist_ok=True)
    img_path = dest_folder / f"{file_stem}_{image_id}.png"
    with open(img_path, "wb") as img_file:
        img_file.write(image_bytes)
    return img_path


def extract_images_from_pdf(pdf_path: Path) -> Dict:
    images_info = []
    doc = fitz.open(pdf_path)
    for page_index, page in enumerate(doc, start=1):
        images = page.get_images(full=True)
        for img_index, img in enumerate(images):
            xref = img[0]
            base_image = doc.extract_image(xref)
            img_bytes = base_image["image"]
            img_path = save_image(img_bytes, OUTPUT_DIR, pdf_path.stem, f"p{page_index}_i{img_index}")
            images_info.append({
                "file_name": pdf_path.name,
                "source_file": str(pdf_path),
                "page": page_index,
                "image_path": str(img_path),
                "type": "pdf"
            })
    return images_info


def extract_images_from_docx(docx_path: Path) -> Dict:
    images_info = []
    doc = docx.Document(docx_path)
    for rel in doc.part._rels:
        rel_obj = doc.part._rels[rel]
        if "image" in rel_obj.target_ref:
            image_data = rel_obj.target_part.blob
            img_path = save_image(image_data, OUTPUT_DIR, docx_path.stem, str(uuid.uuid4())[:8])
            images_info.append({
                "file_name": docx_path.name,
                "source_file": str(docx_path),
                "image_path": str(img_path),
                "type": "docx"
            })
    return images_info


def extract_images_from_pptx(pptx_path: Path) -> Dict:
    images_info = []
    pres = pptx.Presentation(pptx_path)
    for i, slide in enumerate(pres.slides):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                image = shape.image
                img_path = save_image(image.blob, OUTPUT_DIR, pptx_path.stem, f"s{i}_{uuid.uuid4().hex[:6]}")
                images_info.append({
                    "file_name": pptx_path.name,
                    "source_file": str(pptx_path),
                    "slide": i + 1,
                    "image_path": str(img_path),
                    "type": "pptx"
                })
    return images_info


def extract_images_from_xlsx(xlsx_path: Path) -> Dict:
    # Optional: rarely used, can be implemented with xlwings or a temp workaround
    return []  # Most diagrams in xlsx should be extracted as diagrams or text separately


EXTRACTORS = {
    ".pdf": extract_images_from_pdf,
    ".docx": extract_images_from_docx,
    ".pptx": extract_images_from_pptx,
    ".xlsx": extract_images_from_xlsx,
}


def extract_all_images(raw_dir: Path = Path("raw_files"), output_dir: Path = OUTPUT_DIR):
    output_dir.mkdir(parents=True, exist_ok=True)
    summary = []
    for file_path in raw_dir.rglob("*"):
        if not file_path.is_file():
            continue

        ext = file_path.suffix.lower()
        extractor = EXTRACTORS.get(ext)
        if not extractor:
            continue

        print(f"[img] Extracting images from {file_path.name}")
        try:
            images_info = extractor(file_path)
            summary.extend(images_info)
        except Exception as e:
            print(f"[!] Failed to extract images from {file_path.name}: {e}")

    # Save metadata JSON
    with open(output_dir / "image_summary.json", "w", encoding="utf-8") as f:
        import json
        json.dump(summary, f, indent=2, ensure_ascii=False)

    print(f"✅ Extracted {len(summary)} total images.")


if __name__ == "__main__":
    extract_all_images()
